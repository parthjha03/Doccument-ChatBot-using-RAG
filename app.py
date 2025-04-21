from flask import Flask, render_template, request, jsonify
import os
import logging
import PyPDF2
import docx
from pptx import Presentation
import re
import json
import requests
import spacy
import nltk
from utils import count_tokens
from typing import List, Dict, Any
from groq import Groq
from pinecone import Pinecone, ServerlessSpec
from langchain_pinecone import PineconeVectorStore
from rank_bm25 import BM25Okapi
from sentence_transformers import CrossEncoder
import numpy as np
from openai import OpenAI
import chardet
from dotenv import load_dotenv
from deep_translator import GoogleTranslator
from langdetect import detect
import speech_recognition as sr

# Load environment variables from .env file
load_dotenv()

app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = 'uploads'

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

nltk.download('punkt', quiet=True)
nlp = spacy.load("en_core_web_sm")
nlp.max_length = 2000000  # Increase the max_length limit

class EmbeddingWrapper:
    def __init__(self, api_key):
        self.api_key = api_key
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://api.deepinfra.com/v1/openai"
        )

    def create_embeddings(self, texts: List[str], batch_size: int = 10) -> List[List[float]]:
        embeddings = []
        for i in range(0, len(texts), batch_size):
            batch = texts[i:i + batch_size]
            response = self.client.embeddings.create(
                model="BAAI/bge-base-en-v1.5",
                input=batch,
                encoding_format="float"
            )
            embeddings.extend([res.embedding for res in response.data])
        
        # Check the dimension of the embeddings
        if embeddings:
            print(f"Embedding dimension: {len(embeddings[0])}")
        
        return embeddings

    def embed_query(self, query: str) -> List[float]:
        response = self.client.embeddings.create(
            model="BAAI/bge-base-en-v1.5",
            input=[query],
            encoding_format="float"
        )
        return response.data[0].embedding

def improved_get_relevant_chunks(user_question: str, docsearch: PineconeVectorStore, chunks: List[str], top_k: int = 3) -> List[str]:
    vector_results = docsearch.similarity_search(user_question, k=top_k * 2)
    vector_chunks = [doc.page_content for doc in vector_results]
    
    bm25 = BM25Okapi([doc.split() for doc in chunks])
    tokenized_query = user_question.split()
    bm25_scores = bm25.get_scores(tokenized_query)
    top_bm25_indices = np.argsort(bm25_scores)[-top_k * 2:][::-1]
    bm25_chunks = [chunks[i] for i in top_bm25_indices]
    
    combined_chunks = list(dict.fromkeys(vector_chunks + bm25_chunks))
    
    cross_encoder = CrossEncoder('cross-encoder/ms-marco-MiniLM-L-6-v2')
    pairs = [(user_question, chunk) for chunk in combined_chunks]
    cross_encoder_scores = cross_encoder.predict(pairs)
    
    sorted_chunks = [chunk for _, chunk in sorted(zip(cross_encoder_scores, combined_chunks), reverse=True)]
    return sorted_chunks[:top_k]

def semantic_chunking(text: str) -> List[str]:
    doc = nlp(text)
    chunks = []
    current_chunk = []
    current_length = 0
    max_chunk_length = 500

    for sent in doc.sents:
        sent_length = len(sent)
        if current_length + sent_length > max_chunk_length and current_chunk:
            chunks.append(" ".join(current_chunk))
            current_chunk = []
            current_length = 0
        current_chunk.append(sent.text)
        current_length += sent_length

    if current_chunk:
        chunks.append(" ".join(current_chunk))

    return chunks

def populate_pinecone_index(file_path, embedding_wrapper, pc, index_name):
    try:
        # Detect the file encoding
        with open(file_path, "rb") as f:
            raw_data = f.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding']

        # Read the file with the detected encoding
        with open(file_path, "r", encoding=encoding, errors='ignore') as f:
            text = f.read()

        chunks = semantic_chunking(text)

        embeddings = embedding_wrapper.create_embeddings(chunks)

        # Ensure embeddings match the index dimension
        if len(embeddings[0]) != 768:
            raise ValueError(f"Embedding dimension {len(embeddings[0])} does not match index dimension 768")

        data_to_upsert = [
            (f"chunk_{i}", embedding, {"text": chunk})
            for i, (chunk, embedding) in enumerate(zip(chunks, embeddings))
        ]

        index = pc.Index(index_name)
        index.upsert(vectors=data_to_upsert)
        logger.info(f"Upserted {len(data_to_upsert)} vectors to Pinecone index")

        vector_store = PineconeVectorStore(index, embedding_wrapper, "text")
        return vector_store, chunks
    except Exception as e:
        logger.error(f"Error populating Pinecone index: {str(e)}")
        raise

def check_index_populated(index):
    stats = index.describe_index_stats()
    return stats.total_vector_count > 0

def clean_text(text):
    text = re.sub(r'\s+', ' ', text).strip()
    text = ''.join(char for char in text if char.isprintable() or char.isspace())
    return text

def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            for page in reader.pages:
                text += page.extract_text() + "\n"
    except Exception as e:
        logger.error(f"Error extracting text from PDF {pdf_path}: {str(e)}")
    return clean_text(text)

def extract_text_from_docx(docx_path):
    try:
        doc = docx.Document(docx_path)
        text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from DOCX {docx_path}: {str(e)}")
        return ""

def extract_text_from_pptx(pptx_path):
    try:
        prs = Presentation(pptx_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from PPTX {pptx_path}: {str(e)}")
        return ""

def extract_text_from_txt(txt_path):
    try:
        with open(txt_path, 'r', encoding='utf-8') as file:
            text = file.read()
        return clean_text(text)
    except Exception as e:
        logger.error(f"Error extracting text from TXT {txt_path}: {str(e)}")
        return ""

def extract_text_from_files(path):
    all_text = ""

    if os.path.isdir(path):
        for filename in os.listdir(path):
            file_path = os.path.join(path, filename)
            logger.info(f"Processing file: {filename}")

            if filename.lower().endswith('.pdf'):
                text = extract_text_from_pdf(file_path)
            elif filename.lower().endswith('.docx'):
                text = extract_text_from_docx(file_path)
            elif filename.lower().endswith('.pptx'):
                text = extract_text_from_pptx(file_path)
            elif filename.lower().endswith('.txt'):
                text = extract_text_from_txt(file_path)
            else:
                logger.warning(f"Unsupported file type: {filename}")
                continue

            if text:
                all_text += f"--- Start of {filename} ---\n"
                all_text += text
                all_text += f"\n--- End of {filename} ---\n\n"
            else:
                logger.warning(f"No text extracted from {filename}")

    elif os.path.isfile(path):
        filename = os.path.basename(path)
        logger.info(f"Processing file: {filename}")

        if filename.lower().endswith('.pdf'):
            text = extract_text_from_pdf(path)
        elif filename.lower().endswith('.docx'):
            text = extract_text_from_docx(path)
        elif filename.lower().endswith('.pptx'):
            text = extract_text_from_pptx(path)
        elif filename.lower().endswith('.txt'):
            text = extract_text_from_txt(path)
        else:
            logger.warning(f"Unsupported file type: {filename}")
            text = None

        if text:
            all_text += f"--- Start of {filename} ---\n"
            all_text += text
            all_text += f"\n--- End of {filename} ---\n\n"
        else:
            logger.warning(f"No text extracted from {filename}")

    else:
        logger.error(f"Invalid path: {path}")

    return all_text

def upsert_vectors(index, vectors, namespace="default"):
    index.upsert(vectors=vectors, namespace=namespace)

def query_vectors(index, vector, top_k=2, namespace="default", filter=None):
    response = index.query(
        namespace=namespace,
        vector=vector,
        top_k=top_k,
        include_values=True,
        include_metadata=True,
        filter=filter
    )
    return response

def rewrite_query(client, model, original_query, rewrite_conversation_history):
    system_prompt = '''
    You are an AI assistant tasked with reformulating user queries to improve retrieval in a RAG system. The RAG system has information about various documents.
    Given the original query and the conversation history, rewrite it to be more specific, detailed, and likely to retrieve relevant information. Do not make up information that is not in the question, although you are free to add details if they were mentioned earlier in the conversation history.
    Consider the context of the conversation when rewriting the query. You are rewriting the queries such that they can be used for semantic search in a RAG system whose information will be passed on to another LLM for response. Keep this in mind. Not every query needs rewriting; use your judgment on when to rewrite and when not to. ONLY give the rewritten query as output.
    '''

    messages = [
        {"role": "system", "content": system_prompt},
        *rewrite_conversation_history,
        {"role": "user", "content": f"Original query: {original_query}\n\nRewritten query:"}
    ]

    response = client.chat.completions.create(
        model=model,
        messages=messages,
        max_tokens=200,
        temperature=0.3
    )

    rewritten_query = response.choices[0].message.content
    
    rewrite_conversation_history.append({"role": "user", "content": original_query})
    rewrite_conversation_history.append({"role": "assistant", "content": rewritten_query})

    if len(rewrite_conversation_history) > 20:
        rewrite_conversation_history = rewrite_conversation_history[-20:]

    return rewritten_query, rewrite_conversation_history

def chatbot_response(client, model, user_question, relevant_chunks, conversation_history):
    system_prompt = '''
    You are an AI assistant tasked with generating responses based on user questions and relevant information chunks. 
    The information chunks are related to various documents. Use the provided chunks to generate a detailed and accurate response to the user's question. 
    Ensure that the response is relevant and informative, and avoid making up information that is not present in the chunks.
    '''

    messages = [
        {"role": "system", "content": system_prompt},
        *conversation_history,
        {"role": "user", "content": f"User question: {user_question}\n\nRelevant chunks: {relevant_chunks}\n\nResponse:"}
    ]

    response = client.chat.completions.create(
        model=model,
        messages=messages,
        max_tokens=8000,
        temperature=0.3
    )

    chatbot_response = response.choices[0].message.content
    
    conversation_history.append({"role": "user", "content": user_question})
    conversation_history.append({"role": "assistant", "content": chatbot_response})

    if len(conversation_history) > 2:
        conversation_history = conversation_history[-2:]

    return chatbot_response

def cached_translation(text, source_lang, target_lang):
    try:
        if source_lang == target_lang:
            return text
        translator = GoogleTranslator(source=source_lang, target=target_lang)
        return translator.translate(text)
    except Exception as e:
        raise Exception(f"Translation error: {str(e)}")

def translate_to_english(text):
    try:
        detected_lang = detect(text)
        if detected_lang != 'en':
            translator = GoogleTranslator(source=detected_lang, target='en')
            text = translator.translate(text)
        return text, detected_lang
    except Exception as e:
        raise Exception(f"Translation error: {str(e)}")

def translate_response(text, target_lang):
    try:
        return cached_translation(text, 'en', target_lang)
    except Exception as e:
        raise Exception(f"Translation error: {str(e)}")

def get_speech_input():
    recognizer = sr.Recognizer()
    try:
        with sr.Microphone() as source:
            print("Listening... Speak now!")
            audio = recognizer.listen(source, timeout=5, phrase_time_limit=15)
            try:
                text = recognizer.recognize_google(audio)
                detected_lang = detect(text)
                if detected_lang != 'en':
                    translator = GoogleTranslator(source=detected_lang, target='en')
                    text = translator.translate(text)
                return text, detected_lang
            except sr.UnknownValueError:
                print("Could not understand audio")
                return None, None
            except sr.RequestError as e:
                print(f"Could not request results: {str(e)}")
                return None, None
    except Exception as e:
        print(f"Error accessing microphone: {str(e)}")
        return None, None

@app.route('/')
def index():
    return render_template('index.html')


@app.route('/start_recording', methods=['POST'])
def start_recording():
    try:
        data = request.get_json()
        target_lang = data.get('language', 'en')
        
        recognizer = sr.Recognizer()
        with sr.Microphone() as source:
            recognizer.adjust_for_ambient_noise(source)
            audio = recognizer.listen(source, timeout=5)
            
        text = recognizer.recognize_google(audio)
        
        # Translate if not in English
        if target_lang != 'en':
            translator = GoogleTranslator(source='en', target=target_lang)
            text = translator.translate(text)
            
        return jsonify({'success': True, 'text': text})
        
    except sr.UnknownValueError:
        return jsonify({'error': 'Could not understand audio'}), 400
    except sr.RequestError as e:
        return jsonify({'error': f'Could not request results: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/stop_recording', methods=['POST'])
def stop_recording():
    try:
        return jsonify({'success': True})
    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/upload', methods=['POST'])
def upload_file():
    global docsearch, chunks
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400
    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400
    if file:
        file_path = os.path.join(app.config['UPLOAD_FOLDER'], file.filename)
        file.save(file_path)
        text = extract_text_from_files(file_path)
        embedding_wrapper = EmbeddingWrapper(api_key=os.getenv("DEEPINFRA_API_KEY"))
        embeddings = embedding_wrapper.create_embeddings([text])
        
        # Populate Pinecone index with the extracted text and embeddings
        if not check_index_populated(index):
            docsearch, chunks = populate_pinecone_index(file_path, embedding_wrapper, pc, pinecone_index_name)
        else:
            docsearch = PineconeVectorStore(index, embedding_wrapper, "text")
            chunks = semantic_chunking(text)
        
        return jsonify({"message": "File uploaded and processed successfully", "text": text, "embeddings": embeddings}), 200

@app.route('/chat', methods=['POST'])
def chat():
    user_question = request.json['message']
    conversation_history = request.json.get('conversation_history', [])
    rewrite_conversation_history = request.json.get('rewrite_conversation_history', [])
    target_lang = request.json.get('target_lang', 'en')

    # Translate user question to English if necessary
    user_question, detected_lang = translate_to_english(user_question)

    # Rewrite the query
    rewritten_query, rewrite_conversation_history = rewrite_query(client, model, user_question, rewrite_conversation_history)
    logger.info(f"Rewritten query: {rewritten_query}")

    # Get relevant chunks using the improved retrieval method
    relevant_chunks = improved_get_relevant_chunks(rewritten_query, docsearch, chunks)
    
    # Generate response using the original question and relevant chunks
    response = chatbot_response(client, model, user_question, relevant_chunks, conversation_history)

    # Translate response back to the target language
    response = translate_response(response, detected_lang)

    # Update conversation history
    conversation_history.append({"role": "user", "content": user_question})
    conversation_history.append({"role": "assistant", "content": response})

    # Limit conversation history to last 15 exchanges (30 messages)
    if len(conversation_history) > 30:
        conversation_history = conversation_history[-30:]

    return jsonify({'response': response, 'conversation_history': conversation_history, 'rewrite_conversation_history': rewrite_conversation_history})

if __name__ == "__main__":
    model = os.getenv("MODEL")
    groq_api_key = os.getenv("GROQ_API_KEY")
    pinecone_api_key = os.getenv("PINECONE_API_KEY")
    pinecone_index_name = os.getenv("INDEX_NAME")
    deepinfra_api_key = os.getenv("DEEPINFRA_API_KEY")

    client = Groq(api_key=groq_api_key)
    pc = Pinecone(api_key=pinecone_api_key, request_timeout=60)
    embedding_wrapper = EmbeddingWrapper(deepinfra_api_key)

    if pinecone_index_name not in pc.list_indexes().names():
        pc.create_index(
            pinecone_index_name,
            dimension=768,  # Set the correct dimension for the embeddings
            metric='cosine',
            spec=ServerlessSpec(cloud='aws', region='us-east-1')
        )
        logger.info(f"Created new Pinecone index: {pinecone_index_name}")

    index = pc.Index(pinecone_index_name)

    app.run(debug=True)
